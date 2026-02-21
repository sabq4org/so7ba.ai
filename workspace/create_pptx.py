from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xD4, 0xA8, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MID_BLUE = RGBColor(0x24, 0x4E, 0x7A)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rtl_frame(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame

def add_rtl_para(tf, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.RIGHT, space_after=Pt(6)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Arial'
    p.alignment = align
    p.space_after = space_after
    # RTL
    from pptx.oxml.ns import qn
    pPr = p._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')
    return p

def gold_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), top, Inches(11.333), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def slide_title_bar(slide, title_text):
    # Dark blue bar at top
    bar = add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
    tf = rtl_frame(slide, Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.9))
    add_rtl_para(tf, title_text, size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Gold accent line
    gold_line(slide, Inches(1.2))

def content_slide(title, bullets, emojis=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, WHITE)
    slide_title_bar(slide, title)
    
    tf = rtl_frame(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    for i, bullet in enumerate(bullets):
        emoji = emojis[i] if emojis and i < len(emojis) else '◆'
        add_rtl_para(tf, f"{emoji}  {bullet}", size=18, color=DARK_TEXT, space_after=Pt(12))
    return slide

# ==================== SLIDE 1: Cover ====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, DARK_BLUE)

# Gold decorative bar top
add_shape_bg(s1, Inches(0), Inches(0), prs.slide_width, Inches(0.15), GOLD)
# Gold decorative bar bottom
add_shape_bg(s1, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), GOLD)

# Side accent
add_shape_bg(s1, Inches(12.5), Inches(0.15), Inches(0.12), Inches(7.2), GOLD)

tf = rtl_frame(s1, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
add_rtl_para(tf, '🔵  منصة سبق الذكية', size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tf2 = rtl_frame(s1, Inches(1), Inches(3.0), Inches(11), Inches(0.8))
add_rtl_para(tf2, 'المنصة الإخبارية المتقدمة بالذكاء الاصطناعي', size=24, color=GOLD, align=PP_ALIGN.CENTER)

tf3 = rtl_frame(s1, Inches(1), Inches(4.2), Inches(11), Inches(0.6))
add_rtl_para(tf3, '━━━━━━━━━━━━━━━━━━━━', size=18, color=GOLD, align=PP_ALIGN.CENTER)

tf4 = rtl_frame(s1, Inches(1), Inches(5.0), Inches(11), Inches(0.6))
add_rtl_para(tf4, 'sabq.org', size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tf5 = rtl_frame(s1, Inches(1), Inches(5.8), Inches(11), Inches(0.6))
add_rtl_para(tf5, '2026', size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==================== SLIDE 2-14 ====================

content_slide('🌐  نظرة عامة', [
    'منصة إخبارية متقدمة تعمل بالذكاء الاصطناعي',
    '3 لغات: العربية، الإنجليزية، الأردية',
    'نظام بيئي متكامل: إنتاج آلي + توصيات ذكية + توزيع فيروسي',
    'تجربة إخبارية شخصية لكل مستخدم',
    'النطاق: sabq.org — تعمل 24/7',
], ['🤖', '🌍', '⚙️', '👤', '🔗'])

content_slide('🔐  نظام المصادقة والصلاحيات', [
    '8 أدوار مختلفة متعددة المستويات',
    'تسجيل دخول عبر: البريد الإلكتروني / جوجل / أبل',
    'صلاحيات محددة لكل دور',
    'تشفير bcrypt (12 جولة)',
    'حماية CSRF + Rate Limiting',
], ['👥', '🔑', '🛡️', '🔒', '⚡'])

content_slide('📝  إدارة المحتوى', [
    'محرر نصوص متقدم يدعم RTL',
    'إنشاء وتحرير وجدولة المقالات',
    'إدارة الوسائط المتعددة (صور + فيديو)',
    'أنواع متعددة: أخبار عاجلة، تقارير، تحليلات، رأي',
    'شريط أخبار عاجلة متحرك بالوقت الفعلي',
], ['✏️', '📅', '🖼️', '📰', '🔴'])

content_slide('🧠  الذكاء الاصطناعي', [
    'تلخيص تلقائي للمقالات',
    'تصنيف واستخراج كلمات مفتاحية',
    'توليد عناوين بديلة ووصف SEO',
    'مولد المحتوى iFox (مقالات كاملة من فكرة)',
    'استوديو صور ذكي + إنفوجرافيك',
    'التنسيق الذكي التلقائي للمحتوى العربي',
    'المحرك: OpenAI GPT-5.1',
], ['📋', '🏷️', '✨', '🦊', '🎨', '📐', '🚀'])

content_slide('🎯  التوصيات الذكية', [
    'تحليل سلوك كل مستخدم',
    'قسم "لك" + "أكمل القراءة"',
    'خوارزميات تراعي: المقالات المقروءة، الفئات المفضلة، وقت القراءة',
    'إشعارات مخصصة (متصفح + بريد + واتساب)',
    'منع التكرار + توقيت ذكي للإرسال',
], ['📊', '📖', '🔄', '🔔', '⏰'])

content_slide('🎙️  النشرات الصوتية والبودكاست', [
    'تكامل مع ElevenLabs (صوت عربي طبيعي)',
    'نشرات صوتية يومية تلقائية',
    'خلاصة RSS متوافقة مع تطبيقات البودكاست',
    'محرر متقدم مع جدولة',
], ['🗣️', '📻', '📡', '🎛️'])

content_slide('🔗  التوزيع والمشاركة', [
    'روابط قصيرة (7 أحرف) محسنة للمشاركة',
    'بيانات Open Graph تلقائية',
    'تتبع النقرات والمشاركات',
    'نشر عبر واتساب (Twilio)',
    'وكيل بريد إلكتروني ذكي → مقالات تلقائية',
], ['🔗', '📊', '👆', '📱', '📧'])

content_slide('💰  مصادر الإيرادات', [
    'متجر خدمات إعلامية (بيانات صحفية، منشورات سوشال)',
    'إعلانات مدمجة (دفع لكل نقرة + ميزانيات يومية)',
    'منصة ناشرين (بيع محتوى بنظام الحزم)',
    'بطاقات Apple Wallet (صحفية + ولاء)',
    'دعم ضريبة القيمة المضافة',
], ['🛒', '📢', '📦', '💳', '🧾'])

content_slide('🏛️  المنصات المتخصصة', [
    '"عمق" — تحليلات معمقة ومتخصصة',
    '"مُقترب" — محتوى غير إخباري (زوايا + مواضيع)',
    'نظام مقالات الرأي (تقديم + مراجعة + موافقة)',
    'بحث متقدم يدعم العربية (تطبيع + ترتيب بالصلة)',
], ['🔬', '📚', '✍️', '🔍'])

content_slide('⚙️  البنية التقنية', [
    'الواجهة: React + TypeScript + Vite + Tailwind + Radix UI',
    'الخادم: Express.js + TypeScript + Passport.js',
    'قاعدة البيانات: PostgreSQL (Neon Serverless) + Drizzle ORM',
    'التخزين: Google Cloud Storage',
    'الجوال: Capacitor (iOS + Android)',
    'الأمان: bcrypt, CSRF, CSP, HSTS, Rate Limiting',
], ['🖥️', '🔧', '🗄️', '☁️', '📱', '🔐'])

content_slide('📈  الأداء والمعايير', [
    'تحميل كسول (Lazy Loading)',
    'تجمع اتصالات DB محسّن',
    'مهام خلفية متدرجة',
    'مراقبة أداء (Response Time + Memory)',
    'معايير WCAG 2.1 AA للوصولية',
    'واجهات API موثقة بـ Swagger',
], ['⚡', '🔄', '📋', '📊', '♿', '📖'])

content_slide('🌟  الفوائد', [
    'للقراء: تجربة شخصية + ملخصات صوتية + إشعارات ذكية + وصولية',
    'للمحررين: أدوات AI تسرّع الإنتاج + جدولة + تتبع إنتاجية',
    'للإدارة: تحليلات شاملة + صلاحيات + تواصل داخلي + تنبيهات فورية',
], ['👤', '✏️', '📊'])

content_slide('🔭  الرؤية', [
    'الذكاء الاصطناعي أداة تمكينية للصحفي (مو بديل)',
    'أتمتة المهام الروتينية = وقت للعمل الإبداعي',
    'التخصيص محور التجربة — كل قارئ فريد',
    'تعددية لغوية تعكس تنوع الجمهور',
    'الجمع بين سرعة العصر الرقمي وعمق الصحافة الرصينة',
], ['🤝', '⏱️', '🎯', '🌍', '⚖️'])

# ==================== SLIDE 15: Closing ====================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s15, DARK_BLUE)
add_shape_bg(s15, Inches(0), Inches(0), prs.slide_width, Inches(0.15), GOLD)
add_shape_bg(s15, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), GOLD)
add_shape_bg(s15, Inches(0.7), Inches(0.15), Inches(0.12), Inches(7.2), GOLD)

tf = rtl_frame(s15, Inches(1), Inches(2.0), Inches(11), Inches(1.2))
add_rtl_para(tf, '🔵  منصة سبق الذكية', size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tf2 = rtl_frame(s15, Inches(1), Inches(3.5), Inches(11), Inches(0.8))
add_rtl_para(tf2, 'مستقبل الإعلام العربي', size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

tf3 = rtl_frame(s15, Inches(1), Inches(4.5), Inches(11), Inches(0.6))
add_rtl_para(tf3, '━━━━━━━━━━━━━━━━━━━━', size=18, color=GOLD, align=PP_ALIGN.CENTER)

tf4 = rtl_frame(s15, Inches(1), Inches(5.3), Inches(11), Inches(0.6))
add_rtl_para(tf4, 'sabq.org', size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = '/home/openclaw/.openclaw/workspace/سبق_الذكية.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
